import os
import pptx
import pandas as pd
from pptx_replace import replace_text
from pptx import Presentation
from pptx.util import Pt
import io
import pickle
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
import os
import mimetypes
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaFileUpload

# Authenticate and build the Drive API service
def authenticate_drive_api():
    SCOPES = ['https://www.googleapis.com/auth/drive']
    creds = None

    # Path to client_secret.json
    credentials_path = r"E:\Moderate_Automation\online_automate\client_secret_1003512617787-qp854vhunr7bl79lvhe59tnmlqpa4a3p.apps.googleusercontent.com.json"

    # Check if token.pickle exists (for saved credentials)
    if os.path.exists('token.pickle'):
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)

    # If no valid credentials, authenticate using client_secret.json
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES)
            creds = flow.run_local_server(port=0)
        # Save credentials for future runs
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)

    return build('drive', 'v3', credentials=creds)

# List files in Google Drive
def list_drive_files(service):
    results = service.files().list(pageSize=100, fields="files(id, name, mimeType)").execute()
    return results.get('files', [])

# Download a specific file
def download_file(service, file_id, file_name, output_folder="downloads"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    request = service.files().get_media(fileId=file_id)
    file_path = os.path.join(output_folder, file_name)
    with io.FileIO(file_path, 'wb') as file:
        downloader = MediaIoBaseDownload(file, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            print(f"Downloading {file_name}: {int(status.progress() * 100)}%")
    print(f"File downloaded: {file_path}")
    return file_path

# Download a folder and its files
def download_folder(service, folder_id, output_folder="downloads"):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    # Query for files in the folder using mimeType filter for non-folder items
    query = f"'{folder_id}' in parents and mimeType != 'application/vnd.google-apps.folder'"
    results = service.files().list(q=query, fields="files(id, name)").execute()
    folder_files = results.get('files', [])

    if not folder_files:
        print(f"No files found in the folder: {folder_id}")
        return None

    # Create folder path where files will be saved
    folder_path = os.path.join(output_folder)  # Directly assign output folder path
    if not os.path.exists(folder_path):
        os.makedirs(folder_path)

    # Download each file in the folder
    for file in folder_files:
        print(f"Downloading file: {file['name']}")
        download_file(service, file['id'], file['name'], folder_path)
    
    return folder_path

# Download files automatically based on predefined variables
def download_drive_files():
    service = authenticate_drive_api()
    files = list_drive_files(service)

    print("Google Drive content loaded.")

    # templates = {
    #     "Diabetes": r"Diabetes_Template.pptx", 
    #     "Obesity": r"Obesity_Template.pptx" 
    # }

    excel_path = r"Risk_status.xlsx" 
    excel_path_Seq_details = r"Sequencing_Details.xlsx"
    pictures_folder = r"pictures"
    output_folder = r"Outputs"

    # Download the templates
    for template_name, template_path in templates.items():
        print(f"Downloading template: {template_name}")
        match = next((file for file in files if file['name'].lower() == template_path.lower()), None)
        if match:
            download_file(service, match['id'], match['name'])
        else:
            print(f"Template {template_name} not found.")

    # Download the Excel files
    excel_files = {
        "Risk_status": excel_path,
        "Sequencing_Details": excel_path_Seq_details
    }

    for excel_name, excel_file in excel_files.items():
        print(f"Downloading Excel file: {excel_name}")
        match = next((file for file in files if file['name'].lower() == excel_file.lower()), None)
        if match:
            download_file(service, match['id'], match['name'])
        else:
            print(f"Excel file {excel_name} not found.")

    # Download the pictures folder
    print(f"Downloading folder: {pictures_folder}")
    match = next((file for file in files if file['name'].lower() == pictures_folder.lower() and file['mimeType'] == 'application/vnd.google-apps.folder'), None)
    if match:
        download_folder(service, match['id'], output_folder="downloads/pictures")
    else:
        print(f"Folder {pictures_folder} not found.")

    # Download the output folder
    print(f"Downloading folder: {output_folder}")
    match = next((file for file in files if file['name'].lower() == output_folder.lower() and file['mimeType'] == 'application/vnd.google-apps.folder'), None)
    if match:
        download_folder(service, match['id'], output_folder="downloads/Outputs")
    else:
        print(f"Folder {output_folder} not found.")

# Paths to templates and files
templates = {
    "Diabetes": r"downloads\Diabetes_Template.pptx",
    "Obesity": r"downloads\Obesity_Template.pptx"
}

excel_path = r"downloads\Risk_status.xlsx"
excel_path_Seq_details = r"downloads\Sequencing_Details.xlsx"
pictures_folder = r"downloads\pictures"
output_folder = r"downloads\Outputs"
input_folder = r"downloads\Outputs"

# Helper function to map numeric severity values to image filenames
def map_numeric_to_severity(value):
    mapping = {1: "optional", 2: "advised", 3: "essential"}
    return mapping.get(value, None)

def get_image_filename(value):
    if pd.isna(value):
        return None  # Skip replacement for NaN values
    if isinstance(value, (int, float)):  # Map numeric values
        value = map_numeric_to_severity(int(value))
    value = value.strip().replace(" ", "_")  # Ensure consistent naming
    return f"{value}.png"

def adjust_bullet_points(slide):
    """
    Adjusts bullet points in the given slide to use the font 'DM Sans' and size 13.
    Only applies to text frames with bullets, does not change headings or other text styles.
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.level > 0:  # Only adjust bullet points
                    for run in paragraph.runs:
                        run.font.name = "DM Sans"
                        run.font.size = Pt(13)

# Image mappings for slides
image_mapping = {
    "Diabetes": {
        8: [  # 9th page
            ("vitamin-B12", 0),
            ("vitamin-B9 Folic Acid", 1),
            ("vitamin-C", 2),
            ("vitamin-B1", 3),
            ("vitamin-B2", 4),
            ("vitamin-B6 Biotin", 5),
            ("vitamin-E", 6),
            ("vitamin-D", 7),
            ("vitamin-K", 8),
        ],
        9: [  # 10th page
            ("Iron", 0),
            ("Magnesium", 1),
            ("Molybdenum", 2),
            ("Phosphate", 3),
            ("Calcium", 4),
            ("Zinc", 5),
        ],
        3: [  # 4th page
            ("Diabetes", 0),
            ("Insulin Resistance", 1),
            ("Exercise Intolerance", 2)
        ],
        4: [  # 5th page
            ("Carbohydrate Intolerance", 0),
            ("Fat Intolerance", 1),
            ("Protein Intolerance", 2)
        ],
    },
    "Obesity": {
        8: [  # 9th page
            ("vitamin-B12", 0),
            ("vitamin-B9 Folic Acid", 1),
            ("vitamin-C", 2),
            ("vitamin-B1", 3),
            ("vitamin-B2", 4),
            ("vitamin-B6 Biotin", 5),
            ("vitamin-E", 6),
            ("vitamin-D", 7),
            ("vitamin-K", 8),
        ],
        9: [  # 10th page
            ("Iron", 0),
            ("Magnesium", 1),
            ("Molybdenum", 2),
            ("Phosphate", 3),
            ("Calcium", 4),
            ("Zinc", 5),
        ],
        3: [  # 4th page
            ("Obesity", 0),
            ("Insulin Resistance", 1),
            ("Exercise Intolerance", 2)
        ],
        4: [  # 5th page
            ("Carbohydrate Intolerance", 0),
            ("Fat Intolerance", 1),
            ("Protein Intolerance", 2)
        ],
    }
}

# Replace images and adjacent text
def replace_images(prs, row, template_type):
    mappings = image_mapping[template_type]
    for slide_idx, columns in mappings.items():
        slide = prs.slides[slide_idx]
        shapes = [shape for shape in slide.shapes if shape.shape_type == 13]  # Picture placeholders

        for col_name, shape_idx in columns:
            value = row.get(col_name, None)
            image_filename = get_image_filename(value)

            if image_filename:
                image_path = os.path.join(pictures_folder, image_filename)
                if os.path.exists(image_path):
                    if shape_idx < len(shapes):  # Ensure the index is valid
                        shape = shapes[shape_idx]

                        # Replace the image using the placeholder's position and size
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height

                        shape.element.getparent().remove(shape.element)  # Remove the placeholder
                        slide.shapes.add_picture(
                            image_path,
                            left,
                            top,
                            width=width,
                            height=height
                        )
                        print(f"Replaced image for {col_name} on slide {slide_idx + 1} with {image_filename}")

                        # # Check for adjacent text boxes on the right and update them
                        # for text_shape in slide.shapes:
                        #     if (
                        #         text_shape.has_text_frame
                        #         and text_shape.left > left  # Ensure the text box is on the right
                        #         and abs(text_shape.top - top) < height  # Vertically aligned
                        #     ):
                        #         if text_shape.text.strip().lower() == "optional":  # Check for "Optional"
                        #             text_shape.text = image_filename.rsplit('.', 1)[0].capitalize()
                        #             print(f"Updated text adjacent to {col_name} with '{text_shape.text}'")
                        #         break
                        
                        # Check for adjacent text boxes on the right and update them
                        for text_shape in slide.shapes:
                            if (
                                text_shape.has_text_frame
                                and text_shape.left > left  # Ensure the text box is on the right
                                and abs(text_shape.top - top) < height  # Vertically aligned
                            ):
                                if text_shape.text.strip().lower() == "optional":  # Check for "Optional"
                                    # Update the text to the filename (without extension)
                                    text_shape.text = image_filename.rsplit('.', 1)[0].capitalize()
                                    print(f"Updated text adjacent to {col_name} with '{text_shape.text}'")

                                    # Set the font to "DM Sans"
                                    for paragraph in text_shape.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            run.font.name = "DM Sans"  # Set font to "DM Sans"
                                            run.font.size = Pt(13)  # Set font size to 13, adjust as needed
                                    break
                else:
                    print(f"Image not found: {image_filename}")
            else:
                print(f"Invalid value for {col_name}, skipping image replacement.")

# Process a single row and update the PPT
def process_ppt(row, row2, template_path, template_type):
    prs = pptx.Presentation(template_path)

    # Replace placeholders with data
    replace_text(prs, "Patient_Name", str(row2["Patient_Name"]))
    replace_text(prs, "Sample_ID", str(row2["Sample_ID"]))
    replace_text(prs, "Patient_Age", str(row2["Patient_Age"]))
    replace_text(prs, "Patient_Gender", str(row2["Patient_Gender"]))
    #replace_text(prs, "Sequence_length", str(row2["Sequence_length"]))
    replace_text(prs, "Mean_Sequencing_depth", str(row2["Mean_Sequencing_depth"]))
    # Format 'Collection_date' to match "24 Oct 2024"
    if pd.isna(row2["Collection_date"]):
        formatted_date = "Unknown Date"  # Provide a default value or leave it blank
    else:
        formatted_date = row2["Collection_date"].strftime("%d %b %Y")
      # Example: 24 Oct 2024
    replace_text(prs, "Collection_date", formatted_date)

        # Format 'Overall_Alignment_rate' and 'Q30_score' to percentages if present
    if row2.get('Overall_Alignment_rate'):
        try:
            formatted_alignment_rate = f"{float(row2['Overall_Alignment_rate']) * 100:.1f}%"  # Example: 93.2%
        except ValueError:
            formatted_alignment_rate = "Invalid alignment rate"  # Handle cases where the value is not a valid float
    else:
        formatted_alignment_rate = None  # Or you can leave it as an empty string or any placeholder

    if row2.get('Q30_score'):
        try:
            formatted_q30_score = f"{float(row2['Q30_score']) * 100:.1f}%"  # Example: 93.2%
        except ValueError:
            formatted_q30_score = "Invalid Q30 score"  # Handle cases where the value is not a valid float
    else:
        formatted_q30_score = None  # Or an empty string or placeholder

    replace_text(prs, "Overall_Alignment_rate", formatted_alignment_rate)
    replace_text(prs, "Q30_score", formatted_q30_score)

    # Replace images and adjacent text
    replace_images(prs, row, template_type)

    # Adjust bullet points on slide 6 (index 5)
    slide_6 = prs.slides[5]
    adjust_bullet_points(slide_6)

    # Save the updated presentation
    output_path = os.path.join(output_folder, f"{row['Sample_ID']}_Report.pptx")
    prs.save(output_path)
    print(f"Saved: {output_path}")

def replace_text_in_slide(slide, placeholder, replacement_text):
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        # Replace text in the run
                        run.text = run.text.replace(placeholder, replacement_text)
                        # Set the font properties
                        run.font.name = "DM Sans"
                        run.font.size = Pt(13)

# Process patient data and generate PPT
def process_patient(template_path, sample_id, patient_data, indication):
    prs = Presentation(template_path)

    # Create a dictionary with the column names and their corresponding values
    patient_info = patient_data.set_index("Sample_ID").to_dict(orient="index").get(sample_id, {})

    # List of placeholders based on indication
    if indication == "Obesity":
        placeholders = ["Obesity", "Insulin Resistance", "Exercise Intolerance", "Carbohydrate Intolerance", "Fat Intolerance", "Protein Intolerance"]
    elif indication == "Diabetes":
        placeholders = ["Diabetes", "Insulin Resistance", "Exercise Intolerance", "Carbohydrate Intolerance", "Fat Intolerance", "Protein Intolerance"]
    else:
        raise ValueError(f"Unknown indication: {indication}")

    # Replace placeholders in the 6th slide (index 5)
    slide = prs.slides[5]

    for placeholder in placeholders:
        # Check if there's a matching column in the patient's data and replace the placeholder with the data value
        if placeholder in patient_info:
            replace_text_in_slide(slide, placeholder, str(patient_info[placeholder]))

    # Add a patient-specific picture if available
    picture_path = os.path.join(pictures_folder, f"{sample_id}.jpg")
    if os.path.exists(picture_path):
        slide.shapes.add_picture(picture_path, left=Pt(50), top=Pt(400), width=Pt(300), height=Pt(200))

    # Save the updated PowerPoint
    output_path = os.path.join(output_folder, f"{sample_id}_{indication}_Report.pptx")

    try:
        prs.save(output_path)
        print(f"Saved updated presentation for {sample_id} ({indication}) to {output_path}")
    except PermissionError:
        print(f"Permission denied when trying to save the file for {sample_id}. Please check the file path.")

def replace_text_with_bold_and_handle_nan(slide, keywords, patient_info):
    """
    Replace placeholders in text and make specific keywords bold.
    Keep the placeholders if their value is NaN or empty.
    Do not modify the font size or style for "DOCTER'S RECOMMENDATIONS" text.
    """
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame is not None:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                updated_text = paragraph.text
                for placeholder, replacement in patient_info.items():
                    if placeholder in updated_text:
                        if pd.isna(replacement) or replacement is None:
                            updated_text = updated_text.replace(placeholder, "")  # Leave placeholder blank
                        else:
                            updated_text = updated_text.replace(placeholder, str(replacement))
                
                # Handle the specific case for "DOCTER'S RECOMMENDATIONS"
                if "DOCTER'S RECOMMENDATIONS" in updated_text:
                    paragraph.clear()  # Clear original text
                    run = paragraph.add_run()
                    run.text = updated_text  # Set the whole text
                    run.font.name = "Bitter"
                    run.font.size = Pt(28.5)
                    # No need to continue, just skip the keyword bold formatting for this paragraph
                    continue  

                # Handle bold formatting for other keywords (case-insensitive comparison)
                words = updated_text.split()
                paragraph.clear()  # Clear original text
                for word in words:
                    run = paragraph.add_run()
                    run.text = word + " "
                    # Case-insensitive check for keyword
                    if any(keyword.lower() == word.lower().strip(",.") for keyword in keywords):
                        run.font.bold = True
                    run.font.name = "DM Sans"
                    run.font.size = Pt(13)  # Adjust font size if needed


def process_patient(pptx_path, sample_id, patient_data, indication):
    """
    Generate a PowerPoint for a single patient based on their recommendations.
    """
    prs = Presentation(pptx_path)

    # Create a dictionary with the column names and their corresponding values
    patient_info = patient_data.set_index("Sample_ID").to_dict(orient="index").get(sample_id, {})

    # Keywords to make bold
    keywords = ["Diabetes", "Obesity", "Insulin Resistance", "Exercise", "Carbohydrate", "Fat", "Protein", "High", "Intolerance",
                "Mild", "Moderate", "Mild to Moderate", "Advised", "Essential", "Optional"]

    # Replace placeholders and bold keywords in the 6th slide (index 5)
    slide = prs.slides[5]
    replace_text_with_bold_and_handle_nan(slide, keywords, patient_info)

    # Add a patient-specific picture if available
    picture_path = os.path.join(pictures_folder, f"{sample_id}.jpg")
    if os.path.exists(picture_path):
        slide.shapes.add_picture(picture_path, left=Pt(50), top=Pt(400), width=Pt(300), height=Pt(200))

    # Save the updated PowerPoint with the same name in the output folder
    try:
        prs.save(pptx_path)  # Save with the same file name as the original file in the input folder
        print(f"Saved updated presentation for {sample_id} ({indication}) to {pptx_path}")
    except PermissionError:
        print(f"Permission denied when trying to save the file: {pptx_path}. Please close any open files and check permissions.")
    except Exception as e:
        print(f"An error occurred while saving the file: {e}")

# Upload function to replace 'Outputs' folder in Google Drive
def upload_outputs_to_drive(output_folder):
    service = authenticate_drive_api()
    
    # Search for an existing folder named 'Outputs' in Google Drive
    query = "mimeType='application/vnd.google-apps.folder' and name='Outputs' and trashed = false"
    results = service.files().list(q=query, fields="files(id, name)").execute()
    files = results.get('files', [])
    
    # If the folder exists, delete it
    if files:
        folder_id = files[0]['id']
        service.files().delete(fileId=folder_id).execute()
        print(f"Deleted existing 'Outputs' folder from Google Drive.")
    
    # Create a new folder named 'Outputs' in Google Drive
    folder_metadata = {'name': 'Outputs', 'mimeType': 'application/vnd.google-apps.folder'}
    folder = service.files().create(body=folder_metadata, fields='id').execute()
    folder_id = folder['id']
    print(f"Created new 'Outputs' folder in Google Drive with ID: {folder_id}")
    
    # Upload files from the local 'Outputs' folder
    for filename in os.listdir(output_folder):
        file_path = os.path.join(output_folder, filename)
        if os.path.isfile(file_path):
            mime_type, _ = mimetypes.guess_type(file_path)
            media = MediaFileUpload(file_path, mimetype=mime_type)
            file_metadata = {'name': filename, 'parents': [folder_id]}
            
            # Upload file to the newly created folder
            file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            print(f"Uploaded: {filename}")

def main():

    # Load the Excel file containing the patient information
    patient_data = pd.read_excel(excel_path)
    seq_details_data = pd.read_excel(excel_path_Seq_details)

    # Iterate over each row in the DataFrame
    for _, row in patient_data.iterrows():
        sample_id = row["Sample_ID"]
        indication = row["Indication"]  # Diabetes or Obesity
        row2 = seq_details_data[seq_details_data["Sample_ID"] == sample_id]
        if row2.empty:
            print(f"Skipping {sample_id} - no matching data in Seq details")
            continue

        # Convert row2 from DataFrame to Series for easier access
        row2 = row2.iloc[0]


        # Select the appropriate template based on indication
        if indication in templates:
            template_path = templates[indication]
            process_ppt(row, row2, template_path, indication)  # Process the PPT for this patient
        else:
            print(f"Skipping {sample_id} - unknown indication: {indication}")

    # Read data from Excel file
    recommendations_df = pd.read_excel(excel_path, sheet_name="Recommendations")

    # Iterate through each patient in the recommendations
    for _, patient_data in recommendations_df.groupby("Sample_ID"):
        sample_id = patient_data["Sample_ID"].iloc[0]
        indication = patient_data["Indication"].iloc[0]

        # Find the matching PowerPoint files in the input folder
        for pptx_file in os.listdir(input_folder):
            if pptx_file.endswith(".pptx"):
                pptx_path = os.path.join(input_folder, pptx_file)

                # Only process the file if it matches the sample_id (ensuring correct PowerPoint file)
                if sample_id in pptx_file:
                    process_patient(pptx_path, sample_id, patient_data, indication)

if __name__ == "__main__":
    downloaded_paths = download_drive_files() 
    main('')
    # The Google Drive folder ID for 'Moderate_Reports'
    drive_folder_id = 'your_google_drive_folder_id'
    upload_outputs_to_drive(output_folder)